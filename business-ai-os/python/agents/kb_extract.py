"""Agent: Extraction de texte depuis des fichiers uploadés vers la Knowledge Base."""
import re
import unicodedata
from pathlib import Path
from models.schemas import KBExtractRequest, KBExtractResponse


async def extract_document(req: KBExtractRequest) -> KBExtractResponse:
    """Extrait le texte d'un fichier et le sauvegarde dans wiki-data/{userId}/knowledge/."""
    file_path = Path(req.file_path)
    if not file_path.exists():
        return KBExtractResponse(success=False, error="Fichier introuvable", page_count=0, text_path="")

    ext = file_path.suffix.lower()
    text = ""
    page_count = 0

    try:
        if ext == ".pdf":
            text, page_count = _extract_pdf(file_path)
        elif ext in (".docx", ".doc"):
            text, page_count = _extract_docx(file_path)
        elif ext in (".pptx", ".ppt"):
            text, page_count = _extract_pptx(file_path)
        elif ext in (".txt", ".md"):
            text = file_path.read_text(encoding="utf-8", errors="ignore")
            page_count = 1
        else:
            return KBExtractResponse(success=False, error=f"Format non supporté: {ext}", page_count=0, text_path="")
    except Exception as e:
        return KBExtractResponse(success=False, error=str(e), page_count=0, text_path="")

    if not text.strip():
        return KBExtractResponse(success=False, error="Aucun texte extrait", page_count=0, text_path="")

    # Sauvegarder dans wiki-data/{userId}/knowledge/{slug}.md
    kb_dir = Path(req.wiki_base_path) / req.user_id / "knowledge"
    kb_dir.mkdir(parents=True, exist_ok=True)

    slug = _slugify(req.doc_name)
    text_path = kb_dir / f"{req.doc_id}_{slug}.md"

    md_content = f"""# {req.doc_name}

> Catégorie : {req.category}  
> Fichier : {req.original_filename}  
> Type : {ext[1:].upper()}  
> Pages : {page_count}

---

{text}
"""
    text_path.write_text(md_content, encoding="utf-8")

    return KBExtractResponse(
        success=True,
        error="",
        page_count=page_count,
        text_path=str(text_path)
    )


def _extract_pdf(path: Path) -> tuple[str, int]:
    import fitz  # pymupdf
    doc = fitz.open(str(path))
    pages = []
    page_count = len(doc)
    for i in range(page_count):
        page = doc.load_page(i)
        text = page.get_text("text")
        if text.strip():
            pages.append(f"### Page {i+1}\n{text.strip()}")
    doc.close()
    return "\n\n".join(pages), page_count


def _extract_docx(path: Path) -> tuple[str, int]:
    from docx import Document
    doc = Document(str(path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs), 1


def _extract_pptx(path: Path) -> tuple[str, int]:
    from pptx import Presentation
    prs = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slides.append(f"### Slide {i+1}\n" + "\n".join(texts))
    return "\n\n".join(slides), len(prs.slides)


def _slugify(text: str) -> str:
    text = unicodedata.normalize("NFKD", text).encode("ascii", "ignore").decode()
    text = re.sub(r"[^\w\s-]", "", text.lower())
    text = re.sub(r"[-\s]+", "-", text).strip("-")
    return text[:60] or "document"
